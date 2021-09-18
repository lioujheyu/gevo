#!/usr/bin/env python3

import subprocess
import random
import csv
import json
import pickle
import pathlib
import sys
import filecmp
import logging
import io
import os
import time
import gc
import shutil
from itertools import cycle
from threading import Thread
from threading import Lock

# Avoid Not-Find-Display problem
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from deap import base
from deap import creator
from deap import tools

import pptx
import psutil

from rich import print
import rich.progress
from rich.progress import Progress
from rich.logging import RichHandler

import pycuda.driver as cuda

from gevo import irind
from gevo.irind import edits_as_key, llvmMutateWrap
from gevo import fuzzycompare

# critical section of multithreading
lock = Lock()

class evolution:
    # Parameters
    cudaPTX = 'gevo.ptx'
    editFitMap = {}

    # Content
    pop = []
    generation = 0
    presentation = pptx.Presentation()

    mutStats = {
        'failDistribution': {'-1': 0},
        'valid':0, 'invalid':0, 'infinite':0,
        'op_success':{
            'c':0, 'r':0, 'i':0, 's':0, 'm':0, 'p':0, 'x':0
        },
        'op_fail':{
            'c':0, 'r':0, 'i':0, 's':0, 'm':0, 'p':0, 'x':0
        },
    }
    evalStats = {'cx':{'pass':[0],'fail':[0]}, 'mut':{'pass':[0],'fail':[0]}, 'epi':{'pass':[0],'fail':[0]}}

    class _testcase:
        def __init__(self, evolution, idx, kernel, bin, verifier):
            self.idx = idx
            self.kernels = kernel
            self.appBinary = bin
            self.verifier = verifier
            self.args = []
            self.golden = []
            self._evolution = evolution

        def evaluate(self):
            # Since golden has been filled up, passing this testcase into resultcompare
            # won't compare anything which is exactly what we want.
            # Evaluate 3 times and get the minimum number
            fitness = [ self._evolution.execNVprofRetrive(self) for i in range(3)]
            self.fitness = (min([ value[0] for value in fitness ]), min([ value[1] for value in fitness ]))
            if None in self.fitness:
                print(self.args)
                raise Exception("Original binary execution error")

            if self.verifier['mode'] == 'file':
                for fname in self.verifier['output']:
                    golden_filename = fname + '.golden' + str(self.idx)
                    os.rename(fname, golden_filename)
                    self.golden.append(golden_filename)

    def __init__(self, kernel, bin, profile, mutop, timeout=30, fitness='time', popsize=128,
                 llvm_src_filename='cuda-device-only-kernel.ll', use_fitness_map=True, combine_positive_epistasis=False,
                 CXPB=0.8, MUPB=0.1, err_rate='0.01', global_seed=None):
        self.CXPB = CXPB
        self.MUPB = MUPB
        self.err_rate = err_rate
        self.kernels = kernel
        self.appBinary = bin
        self.timeout = timeout
        self.fitness_function = fitness
        self.use_fitness_map = use_fitness_map
        self.combine_positive_epistasis = combine_positive_epistasis
        self.popsize = popsize
        self.mutop = mutop.split(',')
        self.global_seed = global_seed

        try:
            with open('editmap.pickle', 'rb') as editFitMapFile:
                self.editFitMap = pickle.load(editFitMapFile)
                print(f'[Initializing GEVO] Previous EditFitMap found. {len(self.editFitMap)} entries loaded')
        except FileNotFoundError:
            pass

        try:
            with open(llvm_src_filename, 'r') as f:
                self.initSrcEnc = f.read().encode()
        except IOError:
            print("File {} does not exist".format(llvm_src_filename))
            exit(1)

        self.verifier = profile['verify']

        # Tools initialization
        # Detect GPU property
        cuda.init()
        # TODO: check if there are multiple GPUs.
        SM_MAJOR, SM_MINOR = cuda.Device(0).compute_capability()
        self.mgpu = 'sm_' + str(SM_MAJOR) + str(SM_MINOR)
        print(f'[Initializing GEVO] GPU compute capability: {self.mgpu}')

        # check Nvidia Profiler exists
        self.nvprof_path = shutil.which('nvprof')
        if self.nvprof_path is None:
            raise Exception('nvprof cannot be found')
        print(f'[Initializing GEVO] nvprof detected: {self.nvprof_path}')

        # Minimize both performance and error
        creator.create("FitnessMin", base.Fitness, weights=(-1.0, -1.0))
        creator.create("Individual", irind.llvmIRrep, fitness=creator.FitnessMin)
        self.history = tools.History()
        self.toolbox = base.Toolbox()
        self.toolbox.register('mutate', self.mutLLVM)
        self.toolbox.register('mate', self.cxOnePointLLVM)
        # self.toolbox.register('select', tools.selDoubleTournament, fitness_size=2, parsimony_size=1.4, fitness_first=True)
        self.toolbox.register('select', tools.selNSGA2)
        self.toolbox.register('individual', creator.Individual, srcEnc=self.initSrcEnc, mgpu=self.mgpu)
        self.toolbox.register('population', tools.initRepeat, list, self.toolbox.individual)
        # Decorate the variation operators
        self.toolbox.decorate("mate", self.history.decorator)
        self.toolbox.decorate("mutate", self.history.decorator)

        self.stats = tools.Statistics(lambda ind: ind.fitness.values)
        self.stats.register("min", min)
        self.stats.register("max", max)

        self.logbook = tools.Logbook()
        self.paretof = tools.ParetoFront()
        self.logbook.header = "gen", "evals", "min", "max"

        # Set up testcase
        self.origin = creator.Individual(self.initSrcEnc, self.mgpu)
        self.origin.ptx(self.cudaPTX)
        arg_array = [[]]
        for i, arg in enumerate(profile['args']):
            if arg.get('bond', None) is None:
                arg_array_next = [ e[:] for e in arg_array for _ in range(len(arg['value']))]
                arg_array = arg_array_next
                for e1, e2 in zip(arg_array, cycle(arg['value'])):
                    e1.append(e2)
            else:
                for e in arg_array:
                    bonded_arg = arg['bond'][0]
                    bonded_idx = profile['args'][bonded_arg]['value'].index(e[bonded_arg])
                    e.append(arg['value'][bonded_idx])

        arg_array = [ [str(e) for e in args ] for args in arg_array ]

        self.testcase = []
        for i in range(len(arg_array)):
            self.testcase.append( self._testcase(self, i, kernel, bin, profile['verify']) )
        with Progress("[Initializing GEVO] Evaluate original program with test cases",
                      "({task.completed} / {task.total})",
                      rich.progress.TimeElapsedColumn()) as progress:
            task = progress.add_task("", total=len(arg_array))
            for tc, arg in zip(self.testcase, arg_array):
                tc.args = arg
                tc.evaluate()
                progress.update(task, advance=1)

        self.ofits = [ tc.fitness[0] for tc in self.testcase]
        self.oerrs = [ tc.fitness[1] for tc in self.testcase]
        self.origin.fitness.values = (sum(self.ofits)/len(self.ofits), max(self.oerrs))
        self.editFitMap[tuple()] = self.origin.fitness.values
        print(f"Average fitness of the original program: ({self.origin.fitness.values[0]:.2f}, {self.origin.fitness.values[1]:.2f})")
        print("Individual test cases:")
        for fit, err in zip(self.ofits, self.oerrs):
            print(f"\t({fit:.2f}, {err:.2f})")
        self.positive_epistasis = {}
        self.negative_epistasis = {}
        self.need_discussion = {}

    def updateSlideFromPlot(self):
        pffits = [ind.fitness.values for ind in self.paretof]
        fits = [ind.fitness.values for ind in self.pop if ind.fitness.values not in pffits]
        plt.gcf().subplots_adjust(bottom=0.15)
        plt.title("Program variant performance - Generation {}".format(self.generation))
        plt.xlabel("Runtime(ms)")
        plt.ylabel("Error(%)")
        err_rate = float(self.err_rate[0:-2]) if self.err_rate[-1] == '|' else float(self.err_rate)
        plt.ylim(ymin=-(float(err_rate)*100/20), ymax=float(err_rate)*100)
        plt.xticks(rotation=45)
        plt.scatter([fit[0]/1000 for fit in fits], [fit[1]*100 for fit in fits],
                    marker='*', label="dominated")
        plt.scatter([pffit[0]/1000 for pffit in pffits], [pffit[1]*100 for pffit in pffits],
                    marker='o', c='red', label="pareto front")
        plt.scatter(self.origin.fitness.values[0]/1000, self.origin.fitness.values[1]*100,
                    marker='x', c='green', label="origin")
        plt.legend()
        buf = io.BytesIO()
        plt.savefig(buf, format='png')
        buf.seek(0)

        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        left = top = pptx.util.Inches(1)
        pic = slide.shapes.add_picture(buf, left, top)
        buf.close()
        self.presentation.save('progress.pptx')
        plt.clf()

    def writeStage(self):
        pathlib.Path('stage').mkdir(exist_ok=True)
        if self.generation == 0:
            stageFileName = "stage/startedits.json"
        else:
            stageFileName = "stage/" + str(self.generation) + ".json"

        with open(stageFileName, 'w') as fp:
            stage = [{'edits': ind.edits, 'fitness': ind.fitness.values} for ind in self.pop]
            json.dump(stage, fp, indent=2)

    def mutLog(self):
        print("gen {}".format(self.generation), file=self.mutLogFile)
        print("Individual mutation opeartion statistic", file=self.mutLogFile)
        print("       |         c         r         i         s         m         p", file=self.mutLogFile)
        print("success|{:>10}{:>10}{:>10}{:>10}{:>10}{:>10}".format(
            self.mutStats['op_success']['c'],
            self.mutStats['op_success']['r'],
            self.mutStats['op_success']['i'],
            self.mutStats['op_success']['s'],
            self.mutStats['op_success']['m'],
            self.mutStats['op_success']['p'],
        ), file=self.mutLogFile)
        print("fail   |{:>10}{:>10}{:>10}{:>10}{:>10}{:>10}".format(
            self.mutStats['op_fail']['c'],
            self.mutStats['op_fail']['r'],
            self.mutStats['op_fail']['i'],
            self.mutStats['op_fail']['s'],
            self.mutStats['op_fail']['m'],
            self.mutStats['op_fail']['p'],
        ), file=self.mutLogFile)
        print("", file=self.mutLogFile)
        self.mutLogFile.flush()

        numMut = ["Num of Mutation"]
        count = ["Count"]
        for key in self.mutStats['failDistribution']:
            numMut.append(int(key))
            count.append(self.mutStats['failDistribution'][key])
        mutDistWrite = csv.writer(self.mutDistFile, quoting=csv.QUOTE_NONNUMERIC)
        mutDistWrite.writerow(["Gen {}".format(self.generation)])
        mutDistWrite.writerow(numMut)
        mutDistWrite.writerow(count)
        self.mutDistFile.flush()

    def resultCompare(self, stdoutStr, testcase):
        err = 0.0
        if self.verifier['mode'] == 'string':
            src = self.verifier['output']
            golden = testcase.golden
            result = False if src.find(golden) == -1 else True
            return result, err
        elif self.verifier['mode'] == 'thundersvm':
            for line in stdoutStr.splitlines():
                if line.find('Cross Accuracy = ') != -1:
                    accuracy = float(line.replace('Cross Accuracy = ',''))
                    err = 1 - accuracy
                    result = False if err > float(self.err_rate) else True
                    return result, err
            return False, err
        elif self.verifier['mode'] == 'caffe2':
            for line in stdoutStr.splitlines():
                if line.find('Accuracy = ') != -1:
                    accuracy = float(line.replace('Accuracy = ',''))
                    err = 1 - accuracy
                    result = False if err > float(self.err_rate) else True
                    return result, err
            return False, err
        elif self.verifier['mode'] == 'file':
            src = self.verifier['output']
            golden = testcase.golden
            result = True
            for cnt, (s, g) in enumerate(zip(src, golden)):
                fuzzy = self.verifier.get('fuzzy', False)
                if isinstance(fuzzy, list):
                    try:
                        fuzzy = fuzzy[cnt]
                    except IndexError:
                        print("Verification Error: Fuzzy mode is a list but does not match the number of output files")
                elif isinstance(fuzzy, bool):
                    pass
                else:
                    raise Exception("Verification Error: fuzzy mode is not a single or a list of the boolean value")

                if fuzzy:
                    rc, maxerr, avgerr = fuzzycompare.file(s, g, self.err_rate)
                    # if rc < 0:
                    #     raise Exception(msg)
                    result = result & (True if rc==0 else False)
                    err = maxerr if maxerr > err else err
                else:
                    try:
                        result = result & filecmp.cmp(s, g)
                    except IOError:
                        print(f"Verification Error: File {src} or {golden} cannot be found")
            return result, err
        else:
            raise Exception("Verification Error: Unknown comparing mode \"{}\" in the json profile".format(
                self.verifier['mode']))

    def identify_positive_epistasis(self, ind, updateEpiTable=True):
        if (len(ind.edits) < 2):
            return
        tmp_edits = list(ind.edits)
        ret_edits = []
        while len(tmp_edits) != 0:
            edit = tmp_edits[0]
            if isinstance(edit, irind.edit):
                mut_fields = edit[0][1].split(',')
                on_mutated_instruction = False
                for mut_field in mut_fields:
                    try:
                        suffix = mut_field.split('.', maxsplit=1)[1]
                        if suffix.find('OP') == -1:
                            on_mutated_instruction = True
                            break
                    except IndexError:
                        pass
                if on_mutated_instruction:
                    tmp_edits.remove(edit)
                    ret_edits.append(edit)
                    continue
            else: 
                # complex edit cannot fail the test, thus won't depends on other. 
                # Only other depend on them
                tmp_edits.remove(edit)
                ret_edits.append(edit)
                continue

            edit_key = edits_as_key([edit])
            if edit_key not in self.editFitMap.keys():
                if updateEpiTable:
                    tmp_ind = creator.Individual(self.initSrcEnc, self.mgpu, [edit])
                    with lock:
                        self.evaluate(tmp_ind, 'epi')
                else:
                    tmp_edits.remove(edit)
                    ret_edits.append(edit)
                    continue

            if None not in self.editFitMap[edit_key]:
                tmp_edits.remove(edit)
                ret_edits.append(edit)
                continue

            # Epistasis found. Check if epistasis map already has the record
            rest = [ e for e in (tmp_edits+ret_edits) if e != edit ]
            with lock:
                if edit in self.positive_epistasis.keys():
                    found = False
                    for dependant in self.positive_epistasis[edit]:
                        if dependant in rest:
                            found = True
                            tmp_edits.remove(edit)
                            if dependant in ret_edits:
                                ret_edits.remove(dependant)
                            else:
                                tmp_edits.remove(dependant)
                            ret_edits.append(tuple([edit, dependant]))
                            break
                    if found:
                        continue

                if updateEpiTable:
                    identified = False
                    for e_rest in rest:
                        if edits_as_key([edit, e_rest]) not in self.editFitMap:
                            tmp_ind = creator.Individual(self.initSrcEnc, self.mgpu, [edit, e_rest])
                            self.evaluate(tmp_ind, 'epi')

                        if None not in self.editFitMap[edits_as_key([edit, e_rest])]:
                            if e_rest in ret_edits:
                                ret_edits.remove(e_rest)
                            else:
                                tmp_edits.remove(e_rest)
                            tmp_edits.remove(edit)
                            ret_edits.append(tuple([edit, e_rest]))
                            if e_rest not in self.positive_epistasis.setdefault(edit, []):
                                self.positive_epistasis[edit].append(e_rest)
                            identified = True
                            break
                    if identified is True:
                        continue
                    else:
                        self.need_discussion.setdefault(edit, []).append(rest)
            ret_edits.append(edit)
            tmp_edits.remove(edit)
        assert(edits_as_key(ret_edits) == ind.key)
        ind.update_edits(ret_edits)

    def mutLLVM(self, individual, rng:random.Random):
        trial = 0
        operations = self.mutop
        while trial < individual.lineSize*2:
            line1 = rng.randint(1, individual.lineSize)
            line2 = rng.randint(1, individual.lineSize)
            while line1 == line2:
                line2 = rng.randint(1, individual.lineSize)

            seed = rng.getrandbits(16)
            op = rng.choice(operations)
            if op == 'p' or op == 'x':
                rc, _, editUID = llvmMutateWrap(individual.srcEnc, op, str('Rand'), str('Rand'), seed)
            else:
                rc, _, editUID = llvmMutateWrap(individual.srcEnc, op, str(line1), str(line2), seed)
            if rc < 0:
                continue

            if editUID in individual.serialized_edits:
                continue

            try:
                test_ind = creator.Individual(self.initSrcEnc, self.mgpu, individual.edits + [editUID])
            except irind.llvmIRrepRuntimeError:
                continue

            with lock:
                trial = trial + 1
                fit = self.evaluate(test_ind, 'mut')

            if None in fit:
                self.mutStats['invalid'] = self.mutStats['invalid'] + 1
                self.mutStats['op_fail'][op] = self.mutStats['op_fail'][op] + 1
                continue
            self.mutStats['valid'] = self.mutStats['valid'] + 1
            self.mutStats['op_success'][op] = self.mutStats['op_success'][op] + 1
            self.mutStats['failDistribution'][str(trial)] = \
                self.mutStats['failDistribution'][str(trial)] + 1 \
                if str(trial) in self.mutStats['failDistribution'] else 1

            if self.combine_positive_epistasis:
                self.identify_positive_epistasis(test_ind, updateEpiTable=False)
            individual.copy_from(test_ind)
            individual.fitness.values = fit

            return individual,

        print("Cannot get mutant to survive in {} trials".format(individual.lineSize*2))
        self.mutStats['failDistribution']['-1'] = self.mutStats['failDistribution']['-1'] + 1
        with lock:
            fit = self.evaluate(individual, 'mut')
            individual.fitness.values = fit
        return individual,

    def cxOnePointLLVM(self, ind1, ind2, rng:random.Random):
        trial = 0
        if ind1 == ind2:
            return ind1, ind2
        
        editSet1 = set(ind1.edits)
        editSet2 = set(ind2.edits)
        # The following sorting it to make the order deterministic
        intersectionEdits = sorted(list(editSet1.intersection(editSet2)))
        symmetricEdits = sorted(list(editSet1.symmetric_difference(editSet2)))
        if len(symmetricEdits) == 1:
            return ind1, ind2 

        while trial < len(ind1.edits) + len(ind2.edits):
            
            rng.shuffle(symmetricEdits)
            point = rng.randint(1, len(symmetricEdits)-1)
            cmd1 = intersectionEdits + symmetricEdits[:point]
            cmd2 = intersectionEdits + symmetricEdits[point:]

            try:
                child1 = creator.Individual(self.initSrcEnc, self.mgpu, cmd1)
                child2 = creator.Individual(self.initSrcEnc, self.mgpu, cmd2)
            except irind.llvmIRrepRuntimeError:
                trial = trial + 1
                continue

            with lock:
                fit1 = self.evaluate(child1, 'cx')
                fit2 = self.evaluate(child2, 'cx')

            trial = trial + 1
            if None in fit1 and None in fit2:
                continue

            if None not in fit1:
                ind1.copy_from(child1)
                ind1.fitness.values = fit1
            if None not in fit2:
                ind2.copy_from(child2)
                ind2.fitness.values = fit2

            return ind1, ind2

        print("Cannot get crossover to survive in {} trials".format(len(ind1.edits) + len(ind2.edits)))
        return ind1, ind2

    def execNVprofRetrive(self, testcase):
        proc = subprocess.Popen([self.nvprof_path,
                                 '--unified-memory-profiling', 'off',
                                #  '--profile-from-start', 'off',
                                 '--profile-child-processes',
                                 '--profile-api-trace', 'none',
                                 '--system-profiling', 'on',
                                 '--csv',
                                 '-u', 'us',
                                 './' + self.appBinary] + testcase.args,
                                stdout=subprocess.PIPE, stderr=subprocess.PIPE)

        try:
            gc.disable()
            all_time = time.perf_counter()
            stdout, stderr = proc.communicate(timeout=self.timeout) # second
            all_time = time.perf_counter() - all_time
            gc.enable()
            retcode = proc.poll()
            # retcode == 9: error is from testing program, not nvprof
            # retcode == 15: Target program receive segmentation fault
            if retcode == 9 or retcode == 15:
                print('x', end='', flush=True)
                return None, None
            # Unknown nvprof error
            if retcode != 0:
                print(stderr.decode(), file=sys.stderr)
                raise Exception('Unknown nvprof error code:{}'.format(retcode))
        except subprocess.TimeoutExpired:
            # Sometimes terminating nvprof will not terminate the underlying cuda program
            # if that program is corrupted. So issue the kill command to those cuda app first
            print('8', end='', flush=True)
            try:
                parent = psutil.Process(proc.pid)
            except psutil.NoSuchProcess:
                return
            children = parent.children(recursive=True)
            for subproc in children:
                subproc.terminate()
            subprocess.run(['killall', self.appBinary],
                           stdout=subprocess.PIPE,
                           stderr=subprocess.PIPE)

            proc.terminate()
            proc.wait()
            self.mutStats['infinite'] = self.mutStats['infinite'] + 1
            return None, None

        program_output = stdout.decode()
        cmpResult, err = self.resultCompare(program_output, testcase)
        if cmpResult is False:
            print('x', end='', flush=True)
            return None, None
        else:
            print('.', end='', flush=True)
            profile_output = stderr.decode()
            csv_list = list(csv.reader(io.StringIO(profile_output), delimiter=','))

            # search for kernel function(s)
            kernel_time = []
            time_percent = []
            energy = None
            # The stats starts after 5th line
            for line in csv_list[5:]:
                if len(line) == 0:
                    continue
                if line[0] == "GPU activities":
                    for name in self.kernels:
                        # 8th column for name of CUDA function call
                        try:
                            if line[7].split('(')[0] == name:
                                # 3rd column for avg execution time
                                kernel_time.append(float(line[2]))
                                time_percent.append(float(line[1]))
                        except:
                            continue
                if line[0] == "Power (mW)":
                    count = int(line[2])
                    avg_power = float(line[3])
                    # The emprical shows that the sampling frequency is around 50Hz
                    energy = count * avg_power / 20

            if len(self.kernels) == len(kernel_time) and energy is not None:
                if self.fitness_function == 'time' or self.fitness_function == 'all_time':
                    # total_kernel_time = sum(kernel_time)*100 / sum(time_percent)
                    return all_time, err
                    # return sum(kernel_time), err
                elif self.fitness_function == 'kernel_time':
                    # total_kernel_time = sum(kernel_time)*100 / sum(time_percent)
                    total_kernel_time = sum(kernel_time)
                    return total_kernel_time, err
                    # return sum(kernel_time), err
                elif self.fitness_function == 'power':
                    return energy, err
            else:
                print("Can not find kernel \"{}\" from nvprof".format(self.kernels), file=sys.stderr)
                return None, None

    def evaluate(self, individual, mode=None):
        if mode == 'cx':
            print('c', end='', flush=True)
        elif mode == 'mut':
            print('m', end='', flush=True)
        elif mode == 'epi':
            print('e', end='', flush=True)

        # first to check whether we can find the same entry in the editFitmap
        editkey = individual.key
        if self.use_fitness_map is True:
            if editkey in self.editFitMap:
                print('r', end='', flush=True)
                return self.editFitMap[editkey]

        # link
        try:
            individual.ptx(self.cudaPTX)
        except:
            self.editFitMap[editkey] = (None, None)
            if mode is not None:
                self.evalStats[mode]['fail'][-1] = self.evalStats[mode]['fail'][-1] + 1 
            return None, None

        with open('gevo.ll', 'w') as f:
            f.write(individual.srcEnc.decode())

        fits = []
        errs = []
        for tc in self.testcase:
            fitness, err = self.execNVprofRetrive(tc)

            for res_file in self.verifier['output']:
                if os.path.exists(res_file):
                    os.remove(res_file)

            if fitness is None or err is None:
                self.editFitMap[editkey] = (None, None)
                if mode is not None:
                    self.evalStats[mode]['fail'][-1] = self.evalStats[mode]['fail'][-1] + 1
                return None, None

            fits.append(fitness)
            errs.append(err)

        max_err = max(errs)
        avg_fitness = sum(fits)/len(fits)
        # record the edits and the corresponding fitness in the map
        self.editFitMap[editkey] = (avg_fitness, max_err)
        if mode is not None:
            self.evalStats[mode]['pass'][-1] = self.evalStats[mode]['pass'][-1] + 1
        return avg_fitness, max_err

    def evolve(self, resumeGen):
        threadPool = []
        if self.global_seed is not None:
            random.seed(self.global_seed)
        mut_rng = random.Random(random.getrandbits(16))
        cx_rng = random.Random(random.getrandbits(16))

        if resumeGen == -1:
            # PopSize must be a multiple by 4 for SelTournamentDOD to function properly
            popSize = self.popsize
            self.mut_local_rng = [ random.Random(random.getrandbits(16)) for i in range(popSize) ]
            self.cx_local_rng = [ random.Random(random.getrandbits(16)) for i in range(popSize) ]
            print("Initialize the population. Size {}".format(popSize))
            self.pop = self.toolbox.population(n=popSize)

            with Progress(auto_refresh=False) as pbar:
                task1 = pbar.add_task("", total=popSize)
                for cnt, ind in enumerate(self.pop):
                    pbar.update(task1, completed=cnt+1, refresh=True,
                                description=f"Initializing Population with 3 mutations({cnt+1}/{popSize}), "\
                                    f"m:(p:{self.evalStats['mut']['pass'][-1]},f:{self.evalStats['mut']['fail'][-1]}), "\
                                    f"c:(p:{self.evalStats['cx']['pass'][-1]},f:{self.evalStats['cx']['fail'][-1]}), "\
                                    f"e:(p:{self.evalStats['epi']['pass'][-1]},f:{self.evalStats['epi']['fail'][-1]})")
                    _fit = (None, None)
                    while None in _fit:
                        _ind1 = creator.Individual(self.initSrcEnc, self.mgpu)
                        _ind2 = creator.Individual(self.initSrcEnc, self.mgpu)
                        _ind3 = creator.Individual(self.initSrcEnc, self.mgpu)
                        self.toolbox.mutate(_ind1, mut_rng)
                        self.toolbox.mutate(_ind2, mut_rng)
                        self.toolbox.mutate(_ind3, mut_rng)
                        _ind = creator.Individual(self.initSrcEnc, self.mgpu, _ind1.edits + _ind2.edits + _ind3.edits)
                        _fit = self.evaluate(_ind)
                    ind.copy_from(_ind)
                    ind.fitness.values = _fit

            self.writeStage()
        else:
            if resumeGen == 0:
                stageFileName = "stage/startedits.json"
            else:
                stageFileName = "stage/" + str(resumeGen) + ".json"

            try:
                stage = json.load(open(stageFileName))
                allEdits = [ irind.encode_edits_from_list(entry['edits']) for entry in stage ]
            except:
                print(f"GEVO Error in loading stage file \"{stageFileName}\"")
                print(sys.exc_info())
                exit(1)

            popSize = len(allEdits)
            self.mut_local_rng = [ random.Random(random.getrandbits(16)) for i in range(popSize) ]
            self.cx_local_rng = [ random.Random(random.getrandbits(16)) for i in range(popSize) ]
            print("Resume the population from {}. Size {}".format(stageFileName, popSize))
            self.pop = self.toolbox.population(n=popSize)
            self.generation = resumeGen

            resultList = [False] * popSize
            for i, (edits, ind) in enumerate(zip(allEdits, self.pop)):
                ind.update_edits(edits)
                threadPool.append(
                    Thread(target=irind.update_from_edits, args=(i, ind, resultList))
                )
                threadPool[-1].start()

            for i, ind in enumerate(self.pop):
                threadPool[i].join()
                if resultList[i] == False:
                    raise Exception(f"Could not reconstruct ind from edits:{ind.edits}")
                fitness = self.evaluate(ind)
                if None in fitness:
                    for edit in ind.edits:
                        print(edit)
                    raise Exception("Encounter invalid individual during reconstruction")

                # if self.combine_positive_epistasis:
                #     self.identify_positive_epistasis(ind)
                ind.fitness.values = fitness

        # This is to assign the crowding distance to the individuals
        # and also to sort the pop with front rank
        self.pop = self.toolbox.select(self.pop, popSize)
        self.history.update(self.pop)
        record = self.stats.compile(self.pop)
        self.paretof.update(self.pop)
        self.logbook.record(gen=0, evals=popSize, **record)
        print("")
        print(self.logbook.stream)

        self.updateSlideFromPlot()
        self.mutLogFile = open('mut_stat.log', 'w')
        self.mutDistFile = open('mut_dist.csv', 'w')
        self.mutLog()
        minExecTime = record["min"][0]
        print("")

        # pffits = [ind.fitness.values for ind in self.paretof]
        # fits = [ind.fitness.values for ind in self.pop if ind not in pffits]
        # plt.scatter([fit[0] for fit in fits], [fit[1] for fit in fits], marker='*')
        # plt.scatter([pffits[0] for fit in fits], [pffits[1] for fit in fits], marker='o', c=red)
        # plt.savefig(str(self.generation) + '.png')

        rapid_mutation = False
        no_fitness_change = -5 # give more generations at start before rapid mutation
        while True:
            offspring = tools.selTournamentDCD(self.pop, popSize)
            # Clone the selected individuals
            offspring = list(map(self.toolbox.clone, offspring))

            paretofGen = tools.sortNondominated(self.pop, popSize, first_front_only=True)
            paretofGen[0].sort(key=lambda ind: ind.fitness.values[0])
            with open("g{}_noerr.ll".format(self.generation), 'w') as f:
                f.write(paretofGen[0][-1].srcEnc.decode())
            with open("g{}_noerr.edit".format(self.generation), 'w') as f:
                print(paretofGen[0][-1].edits, file=f)
            with open("g{}_maxerr.ll".format(self.generation), 'w') as f:
                f.write(paretofGen[0][0].srcEnc.decode())
            with open("g{}_maxerr.edit".format(self.generation), 'w') as f:
                print(paretofGen[0][0].edits, file=f)
            with open("editmap.pickle", 'wb') as emfile:
                pickle.dump(self.editFitMap, emfile)

            self.generation = self.generation + 1
            for _, value in self.evalStats.items():
                value['pass'].append(0)
                value['fail'].append(0)

            threadPool.clear()
            for cnt, (child1, child2) in enumerate(zip(offspring[::2], offspring[1::2])):
                if len(child1.edits) < 2 and len(child2.edits) < 2:
                    continue
                if cx_rng.random() < self.CXPB:
                    threadPool.append(
                        Thread(target=self.toolbox.mate, args=(child1, child2, self.cx_local_rng[cnt])))
                    threadPool[-1].start()
            for thread in threadPool:
                thread.join()

            threadPool.clear()
            for cnt, mutant in enumerate(offspring):
                if mut_rng.random() < self.MUPB:
                    del mutant.fitness.values
                    threadPool.append(
                        Thread(target=self.toolbox.mutate, args=(mutant, self.mut_local_rng[cnt])))
                    threadPool[-1].start()
            for thread in threadPool:
                thread.join()

            dead_inds = [ ind for ind in self.pop if None in ind.fitness.values ]
            assert(len(dead_inds) == 0)

            elite = self.toolbox.select(self.pop, int(popSize/64))
            if self.combine_positive_epistasis:
                for ind in elite:
                    self.identify_positive_epistasis(ind)

            self.pop = self.toolbox.select(elite + offspring, popSize)
            record = self.stats.compile(self.pop)
            self.logbook.record(gen=self.generation, evals=popSize, **record)
            self.paretof.update(self.pop)

            print("")
            print(f"m:(p:{self.evalStats['mut']['pass'][-1]},f:{self.evalStats['mut']['fail'][-1]}), "\
                  f"c:(p:{self.evalStats['cx']['pass'][-1]},f:{self.evalStats['cx']['fail'][-1]}), "\
                  f"e:(p:{self.evalStats['epi']['pass'][-1]},f:{self.evalStats['epi']['fail'][-1]})")
            print(self.logbook.stream)
            self.mutLog()
            self.updateSlideFromPlot()
            self.writeStage()
            print("") # an empty line as a generation separator
