#!/usr/bin/env python3

import argparse
import subprocess
import random
import csv
import json
import pathlib
import sys
import filecmp
import io
import os
from itertools import cycle
from threading import Thread
from threading import Lock

# too avoid running into Not-Find-Display problem
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
# import networkx as nx
import numpy
from deap import base
from deap import creator
from deap import tools
import pptx

sys.path.append('/home/jliou4/genetic-programming/cuda_evolve')
import irind
from irind import llvmMutateWrap
from irind import update_from_edits
import fuzzycompare

# critical section of multithreading
lock = Lock()

class evolution:
    # Parameters
    log = open('debug_log', 'w')
    cudaPTX = 'a.ptx'

    # Content
    pop = []
    generation = 0
    presentation = pptx.Presentation()

    mutStats = {
        'valid':0, 'invalid':0, 'infinite':0,
        'maxFit':[], 'avgFit':[], 'minFit':[]
    }

    class _testcase:
        def __init__(self, evolution, idx, kernel, bin, verifier):
            self.idx = idx
            self.kernels = kernel
            self.appBinary = bin
            self.verifier = verifier
            self.args = []
            self.golden = []
            self._evolution = evolution

        def evaluate(self):
            # Since golden has been filled up, passing this testcase into resultcompare
            # won't compare anything which is exactly what we want.
            self.fitness = self._evolution.execNVprofRetrive(self)
            if None in self.fitness:
                print(self.args)
                raise Exception("Original binary execution error")

            for fname in self.verifier['output']:
                golden_filename = fname + '.golden' + str(self.idx)
                os.rename(fname, golden_filename)
                self.golden.append(golden_filename)

    def __init__(self, kernel, bin, profile, timeout=30, fitness='time',
                 llvm_src_filename='cuda-device-only-kernel.ll',
                 CXPB=0.8, MUPB=0.1):
        self.CXPB = CXPB
        self.MUPB = MUPB
        self.kernels = kernel
        self.appBinary = bin
        # self.appArgs = "" if args is None else args
        self.timeout = timeout
        self.fitness_function = fitness

        try:
            with open(llvm_src_filename, 'r') as f:
                self.initSrcEnc = f.read().encode()
        except IOError:
            print("File {} does not exist".format(llvm_src_filename))
            exit(1)

        self.verifier = profile['verify']

        # tools initialization
        # Minimize both performance and error
        creator.create("FitnessMin", base.Fitness, weights=(-1.0, -1.0))
        creator.create("Individual", irind.llvmIRrep, fitness=creator.FitnessMin)
        self.history = tools.History()
        self.toolbox = base.Toolbox()
        self.toolbox.register('mutate', self.mutLLVM)
        self.toolbox.register('mate', self.cxOnePointLLVM)
        # self.toolbox.register('select', tools.selDoubleTournament, fitness_size=2, parsimony_size=1.4, fitness_first=True)
        self.toolbox.register('select', tools.selNSGA2)
        self.toolbox.register('individual', creator.Individual, srcEnc=self.initSrcEnc)
        self.toolbox.register('population', tools.initRepeat, list, self.toolbox.individual)
        # Decorate the variation operators
        self.toolbox.decorate("mate", self.history.decorator)
        self.toolbox.decorate("mutate", self.history.decorator)

        self.stats = tools.Statistics(lambda ind: ind.fitness.values)
        self.stats.register("min", min)
        self.stats.register("max", max)

        self.logbook = tools.Logbook()
        self.paretof = tools.ParetoFront()
        self.logbook.header = "gen", "evals", "min", "max"

        # Set up testcase
        self.origin = creator.Individual(self.initSrcEnc)
        self.origin.ptx(self.cudaPTX)
        arg_array = [[]]
        for i, arg in enumerate(profile['args']):
            if arg.get('bond', None) is None:
                arg_array_next = [ e[:] for e in arg_array for _ in range(len(arg['value']))]
                arg_array = arg_array_next
                for e1, e2 in zip(arg_array, cycle(arg['value'])):
                    e1.append(e2)
            else:
                for e in arg_array:
                    bonded_arg = arg['bond'][0]
                    bonded_idx = profile['args'][bonded_arg]['value'].index(e[bonded_arg])
                    e.append(arg['value'][bonded_idx])

        arg_array = [ [str(e) for e in args ] for args in arg_array ]

        self.testcase = []
        for i in range(len(arg_array)):
            self.testcase.append( self._testcase(self, i, kernel, bin, profile['verify']) )
        print("evalute testcase as golden..", end='', flush=True)
        for i, (tc, arg) in enumerate(zip(self.testcase, arg_array)):
            tc.args = arg
            print("{}..".format(i+1), end='', flush=True)
            tc.evaluate()
        print("done", flush=True)

        fits = [ tc.fitness[0] for tc in self.testcase]
        errs = [ tc.fitness[1] for tc in self.testcase]
        self.origin.fitness.values = (sum(fits)/len(fits), max(errs))

    def updateSlideFromPlot(self):
        pffits = [ind.fitness.values for ind in self.paretof]
        fits = [ind.fitness.values for ind in self.pop if ind not in pffits]
        plt.gcf().subplots_adjust(bottom=0.15)
        plt.title("Program variant performance - Generation {}".format(self.generation))
        plt.xlabel("Runtime(ms)")
        plt.ylabel("Error(%)")
        plt.ylim(ymin=-0.1, ymax=1.0)
        plt.xticks(rotation=45)
        plt.scatter([fit[0]/1000 for fit in fits], [fit[1]*100 for fit in fits],
                    marker='*', label="dominated")
        plt.scatter([pffit[0]/1000 for pffit in pffits], [pffit[1]*100 for pffit in pffits],
                    marker='o', c='red', label="pareto front")
        plt.scatter(self.origin.fitness.values[0]/1000, self.origin.fitness.values[1]*100,
                    marker='x', c='green', label="origin")
        plt.legend()
        buf = io.BytesIO()
        plt.savefig(buf, format='png')
        buf.seek(0)

        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        left = top = pptx.util.Inches(1)
        pic = slide.shapes.add_picture(buf, left, top)
        buf.close()
        self.presentation.save('progress.pptx')
        plt.clf()

    def writeStage(self):
        pathlib.Path('stage').mkdir(exist_ok=True)
        if self.generation == 0:
            stageFileName = "stage/startedits.json"
        else:
            stageFileName = "stage/" + str(self.generation) + ".json"

        with open(stageFileName, 'w') as fp:
            count = 0
            allEdits = [ind.edits for ind in self.pop]
            json.dump(allEdits, fp, indent=2)

    def resultCompare(self, stdoutStr, testcase):
        src = self.verifier['output']
        golden = testcase.golden

        err = 0.0
        if self.verifier['mode'] == 'string':
            result = False if src.find(golden) == -1 else True
            return result, err
        elif self.verifier['mode'] == 'file':
            result = True
            for s, g in zip(src, golden):
                if self.verifier.get('fuzzy', False) == False:
                    try:
                        result = result & filecmp.cmp(s, g)
                    except IOError:
                        print("File {} or {} cannot be found".format(src, golden))
                else:
                    rc, msg, maxerr, avgerr = fuzzycompare.file(s, g)
                    if rc < 0:
                        raise Exception(msg)
                    result = result & (True if rc==0 else False)
                    err = maxerr if maxerr > err else err
            return result, err
        else:
            raise Exception("Unknown comparing mode \"{}\" from compare.json".format(
                self.verifier['mode']))

    def mutLLVM(self, individual):
        trial = 0
        # cut, replace, insert, swap, move, operand replace
        operations = ['c', 'r', 'i', 's', 'm', 'p']
        while trial < individual.lineSize:
            line1 = random.randint(1, individual.lineSize)
            line2 = random.randint(1, individual.lineSize)
            while line1 == line2:
                line2 = random.randint(1, individual.lineSize)

            op = random.choice(operations)
            if op == 'p':
                rc, mutateSrc, editUID = llvmMutateWrap(individual.srcEnc, op, str('Rand'), str('Rand'))
            else:
                rc, mutateSrc, editUID = llvmMutateWrap(individual.srcEnc, op, str(line1), str(line2))
            if rc < 0:
                continue

            test_ind = creator.Individual(self.initSrcEnc)
            test_ind.edits[:] = individual.edits + [editUID]
            test_ind.rearrage()
            if test_ind.update_from_edits() == False:
                continue

            with lock:
                trial = trial + 1
                fit = self.evaluate(test_ind)
                print('m', end='', flush=True)
            if None in fit:
                continue

            individual.update(srcEnc=test_ind.srcEnc)
            individual.edits.append(editUID)
            individual.rearrage()
            individual.fitness.values = fit
            return individual,

        print("Cannot get mutant to be compiled in {} trials".format(individual.lineSize))
        return individual,

    def cxOnePointLLVM(self, ind1, ind2):
        shuffleEdits = ind1.edits + ind2.edits
        random.shuffle(shuffleEdits)
        point = random.randint(1, len(shuffleEdits)-1)
        cmd1 = shuffleEdits[:point]
        cmd2 = shuffleEdits[point:]
        cmd1 = irind.rearrage(cmd1)
        cmd2 = irind.rearrage(cmd2)

        child1 = creator.Individual(self.initSrcEnc)
        child1.edits = list(cmd1)
        child1.update_from_edits(sweepEdits=False)
        child2 = creator.Individual(self.initSrcEnc)
        child2.edits = list(cmd2)
        child2.update_from_edits(sweepEdits=False)

        with lock:
            fit1 = self.evaluate(child1)
            fit2 = self.evaluate(child2)
            print('c', end='', flush=True)

        if None not in fit1:
            ind1 = child1
            ind1.fitness.values = fit1
        if None not in fit2:
            ind2 = child2
            ind2.fitness.values = fit2

        return ind1, ind2

    def execNVprofRetrive(self, testcase):
        proc = subprocess.Popen(['/usr/local/cuda/bin/nvprof',
                                 '--unified-memory-profiling', 'off',
                                 '--profile-from-start', 'off',
                                 '--system-profiling', 'on',
                                 '--csv',
                                 '-u', 'us',
                                 './' + self.appBinary] + testcase.args,
                                stdout=subprocess.PIPE, stderr=subprocess.PIPE)

        try:
            stdout, stderr = proc.communicate(timeout=self.timeout) # second
            retcode = proc.poll()
            # retcode == 9: error is from testing program, not nvprof
            # retcode == 15: Target program receive segmentation fault
            if retcode == 9 or retcode == 15:
                print('x', end='', flush=True)
                self.mutStats['invalid'] = self.mutStats['invalid'] + 1
                return None, None
            # Unknown nvprof error
            if retcode != 0:
                print(stderr.decode(), file=sys.stderr)
                raise Exception('nvprof error')
        except subprocess.TimeoutExpired:
            # Sometimes terminating nvprof will not terminate the underlying cuda program
            # if that program is corrupted. So issue the kill command to those cuda app first
            print('8', end='', flush=True)
            subprocess.run(['killall', self.appBinary],
                           stdout=subprocess.PIPE,
                           stderr=subprocess.PIPE)
            proc.kill()
            proc.wait()
            self.mutStats['infinite'] = self.mutStats['infinite'] + 1
            return None, None

        program_output = stdout.decode()
        cmpResult, err = self.resultCompare(program_output, testcase)
        if cmpResult is False:
            print('x', end='', flush=True)
            self.mutStats['invalid'] = self.mutStats['invalid'] + 1
            return None, None
        else:
            print('.', end='', flush=True)
            self.mutStats['valid'] = self.mutStats['valid'] + 1
            profile_output = stderr.decode()
            csv_list = list(csv.reader(io.StringIO(profile_output), delimiter=','))

            # search for kernel function(s)
            kernel_time = []
            energy = None
            # The stats starts after 5th line
            for line in csv_list[5:]:
                if len(line) == 0:
                    continue
                if line[0] == "GPU activities":
                    for name in self.kernels:
                        # 8th column for name of CUDA function call
                        try:
                            if line[7].split('(')[0] == name:
                                # 3rd column for avg execution time
                                kernel_time.append(float(line[2]))
                        except:
                            continue
                if line[0] == "Power (mW)":
                    count = int(line[2])
                    avg_power = float(line[3])
                    # The emprical shows that the sampling frequency is around 50Hz
                    energy = count * avg_power / 20

            if len(self.kernels) == len(kernel_time) and energy is not None:
                if self.fitness_function == 'time':
                    return sum(kernel_time), err
                elif self.fitness_function == 'power':
                    return energy, err
            else:
                print("Can not find kernel \"{}\" from nvprof".format(self.kernels), file=sys.stderr)
                return None, None

    def evaluate(self, individual):
        # link
        try:
            individual.ptx(self.cudaPTX)
        except:
            return None, None

        with open('a.ll', 'w') as f:
            f.write(individual.srcEnc.decode())

        fits = []
        errs = []
        for tc in self.testcase:
            fitness, err = self.execNVprofRetrive(tc)
            if fitness is None or err is None:
                return None, None

            fits.append(fitness)
            errs.append(err)

        max_err = max(errs)
        avg_fitness = sum(fits)/len(fits)
        return avg_fitness, max_err

    def evolve(self, resumeGen):
        threadPool = []
        if resumeGen == -1:
            # PopSize must be a multiple by 4 for SelTournamentDOD to function properly
            popSize = 100
            print("Initialize the population. Size {}".format(popSize))
            self.pop = self.toolbox.population(n=popSize)

            # Initial 3x mutate to get diverse population
            for ind in self.pop:
                self.toolbox.mutate(ind)
                self.toolbox.mutate(ind)
                self.toolbox.mutate(ind)
            self.writeStage()
        else:
            if resumeGen == 0:
                stageFileName = "stage/startedits.json"
            else:
                stageFileName = "stage/" + str(resumeGen) + ".json"

            try:
                allEdits = json.load(open(stageFileName))
            except:
                print(sys.exc_info())
                exit(1)

            popSize = len(allEdits)
            # popSize = 10
            print("Resume the population from {}. Size {}".format(stageFileName, popSize))
            self.pop = self.toolbox.population(n=popSize)
            self.generation = resumeGen

            resultList = [False] * popSize
            for i, (edits, ind) in enumerate(zip(allEdits, self.pop)):
                editsList = []
                for editG in edits:
                    editsList.append([(e[0], e[1]) for e in editG])
                ind.edits = editsList
                threadPool.append(
                    Thread(target=update_from_edits, args=(i, ind, resultList))
                )
                threadPool[-1].start()

            for i, ind in enumerate(self.pop):
                threadPool[i].join()
                if resultList[i] == False:
                    raise Exception("Could not reconstruct ind from edits:{}".format(ind.edits))
                fitness = self.evaluate(ind)
                if None in fitness:
                    for edit in ind.edits:
                        print(edit)
                    raise Exception("Encounter invalid individual during reconstruction")
                ind.fitness.values = fitness

        # This is to assign the crowding distance to the individuals
        # and also to sort the pop with front rank
        self.pop = self.toolbox.select(self.pop, popSize)
        self.history.update(self.pop)
        record = self.stats.compile(self.pop)
        self.paretof.update(self.pop)
        self.logbook.record(gen=0, evals=popSize, **record)
        print("")
        print(self.logbook.stream)
        self.updateSlideFromPlot()

        # pffits = [ind.fitness.values for ind in self.paretof]
        # fits = [ind.fitness.values for ind in self.pop if ind not in pffits]
        # plt.scatter([fit[0] for fit in fits], [fit[1] for fit in fits], marker='*')
        # plt.scatter([pffits[0] for fit in fits], [pffits[1] for fit in fits], marker='o', c=red)
        # plt.savefig(str(self.generation) + '.png')

        while True:
            offspring = tools.selTournamentDCD(self.pop, popSize)
            # Clone the selected individuals
            offspring = list(map(self.toolbox.clone, offspring))

            # for i, ind in enumerate(self.paretof):
            paretofGen = tools.sortNondominated(self.pop, popSize, first_front_only=True)
            for i, ind in enumerate(paretofGen[0]):
                with open("g{}_pf{}.ll".format(self.generation, i), 'w') as f:
                    f.write(ind.srcEnc.decode())
                with open("g{}_pf{}.edit".format(self.generation, i), 'w') as f:
                    print(ind.edits, file=f)

            self.generation = self.generation + 1

            threadPool.clear()
            for child1, child2 in zip(offspring[::2], offspring[1::2]):
                if len(child1) < 2 and len(child2) < 2:
                    continue
                if random.random() < self.CXPB:
                    threadPool.append(
                        Thread(target=self.toolbox.mate, args=(child1, child2)))
                    threadPool[-1].start()
            for thread in threadPool:
                thread.join()

            threadPool.clear()
            for mutant in offspring:
                if random.random() < self.MUPB:
                    del mutant.fitness.values
                    threadPool.append(
                        Thread(target=self.toolbox.mutate, args=(mutant,)))
                    threadPool[-1].start()
            for thread in threadPool:
                thread.join()


            # Evaluate the individuals with an invalid fitness
            invalid_ind = [ind for ind in offspring if not ind.fitness.valid]
            fitnesses = [self.evaluate(ind) for ind in invalid_ind]
            for ind, fit in zip(invalid_ind, fitnesses):
                ind.fitness.values = fit

            self.pop = self.toolbox.select(self.pop + offspring, popSize)
            record = self.stats.compile(self.pop)
            self.logbook.record(gen=self.generation, evals=popSize, **record)
            self.paretof.update(self.pop)

            print("")
            print(self.logbook.stream)
            self.updateSlideFromPlot()
            self.writeStage()

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="Evolve CUDA kernel function")
    parser.add_argument('-P', '--profile_file', type=str, required=True,
        help="Specify the profile file that contains all application execution and testing information")
    # parser.add_argument('-k', '--kernel', type=str,
    #     help="Target kernel function of the given CUDA application. Use comma to separate kernels.")
    parser.add_argument('-r', '--resume', type=int, default=-1,
        help="Resume the process from genetating the population by reading stage/<RESUME>.json")
    parser.add_argument('-t', '--timeout', type=int, default=30,
        help="The timeout period to evaluate the CUDA application")
    parser.add_argument('-fitf', '--fitness_function', type=str, default='time',
        help="What is the target fitness for the evolution. Default ot execution time. Can be changed to power")
    # parser.add_argument('binary',help="Binary of the CUDA application", nargs='?', default='a.out')
    # parser.add_argument('args',help="arguments for the application binary", nargs=argparse.REMAINDER)
    args = parser.parse_args()

    try:
        profile = json.load(open(args.profile_file))
    except:
        print(sys.exc_info())
        exit(-1)

    evo = evolution(
        kernel=profile['kernels'],
        bin=profile['binary'],
        profile=profile,
        timeout=args.timeout,
        fitness=args.fitness_function)

    print("      Target CUDA program: {}".format(profile['binary']))
    print("Args for the CUDA program:")
    for tc in evo.testcase:
        print("\t{}".format(" ".join(tc.args)))
    print("           Target kernels: {}".format(" ".join(profile['kernels'])))
    print("       Evaluation Timeout: {}".format(args.timeout))
    print("         Fitness function: {}".format(args.fitness_function))

    try:
        evo.evolve(args.resume)
    except KeyboardInterrupt:
        subprocess.run(['killall', args.binary])
        print("valid variant:   {}".format(evo.stats['valid']))
        print("invalid variant: {}".format(evo.stats['invalid']))
        print("infinite variant:{}".format(evo.stats['infinite']))
        if evo.generation > 0:
            evo.presentation.save('progress.pptx')
